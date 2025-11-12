import os
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import pytesseract
from PIL import Image

class ReaderAgent:
    def extract_content(self, file_path):
        """
        Extracts text content from a file based on its extension.
        """
        _, file_extension = os.path.splitext(file_path)
        file_extension = file_extension.lower()

        if file_extension == ".pdf":
            return self._extract_text_from_pdf(file_path)
        elif file_extension == ".pptx":
            return self._extract_text_from_pptx(file_path)
        elif file_extension == ".docx":
            return self._extract_text_from_docx(file_path)
        elif file_extension in [".png", ".jpg", ".jpeg"]:
            return self._extract_text_from_image(file_path)
        else:
            print(f"Unsupported file type: {file_extension}")
            return ""

    def _extract_text_from_pdf(self, file_path):
        """Extracts text from a PDF file."""
        try:
            doc = fitz.open(file_path)
            text = []
            for page in doc:
                text.append(page.get_text())
            return "\n".join(text)
        except Exception as e:
            print(f"Error extracting from .pdf: {e}")
            return ""

    def _extract_text_from_pptx(self, file_path):
        """Extracts text from a PowerPoint file."""
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"Error extracting from .pptx: {e}")
            return ""

    def _extract_text_from_docx(self, file_path):
        """Extracts text from a Word document."""
        try:
            doc = Document(file_path)
            text = [p.text for p in doc.paragraphs]
            return "\n".join(text)
        except Exception as e:
            print(f"Error extracting from .docx: {e}")
            return ""

    def _extract_text_from_image(self, file_path):
        """Extracts text from an image using OCR."""
        try:
            return pytesseract.image_to_string(Image.open(file_path))
        except Exception as e:
            print(f"Error extracting from image: {e}")
            return ""
