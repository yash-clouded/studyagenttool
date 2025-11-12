import fitz  # PyMuPDF
import os
from pathlib import Path

def extract_text_from_pdf(path: str) -> str:
    """Extract text from PDF files."""
    doc = fitz.open(path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_pptx(path: str) -> str:
    """Extract text from PowerPoint slides (.pptx files)."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("python-pptx is required for PowerPoint support. Install with: pip install python-pptx")
    
    prs = Presentation(path)
    text = ""
    for slide_num, slide in enumerate(prs.slides, 1):
        text += f"\n--- Slide {slide_num} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_image(path: str) -> str:
    """Extract text from images (handwritten notes, screenshots) using OCR."""
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        raise ImportError("pytesseract and Pillow are required for image OCR. Install with: pip install pytesseract Pillow")
    
    try:
        img = Image.open(path)
        text = pytesseract.image_to_string(img)
        return text
    except Exception as e:
        raise Exception(f"Failed to extract text from image {path}: {e}")

def extract_text_from_docx(path: str) -> str:
    """Extract text from Word documents (.docx files)."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("python-docx is required for Word document support. Install with: pip install python-docx")
    
    doc = Document(path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_txt(path: str) -> str:
    """Extract text from plain text files."""
    with open(path, 'r', encoding='utf-8') as f:
        return f.read()

def extract_text_from_file(path: str) -> str:
    """
    Extract text from various file formats.
    Supports: PDF, PPTX, DOCX, TXT, PNG, JPG, JPEG, GIF, BMP
    """
    file_ext = Path(path).suffix.lower()
    
    if file_ext == '.pdf':
        return extract_text_from_pdf(path)
    elif file_ext == '.pptx':
        return extract_text_from_pptx(path)
    elif file_ext == '.docx':
        return extract_text_from_docx(path)
    elif file_ext == '.txt':
        return extract_text_from_txt(path)
    elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
        return extract_text_from_image(path)
    else:
        raise ValueError(f"Unsupported file format: {file_ext}. Supported formats: PDF, PPTX, DOCX, TXT, PNG, JPG, JPEG, GIF, BMP")
